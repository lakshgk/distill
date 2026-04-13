"""
distill.parsers.pptx
~~~~~~~~~~~~~~~~~~~~
Parser for Microsoft PowerPoint presentations (.pptx).

Primary path:   python-pptx (structural slide extraction)
Legacy stub:    .ppt raises ParseError with LibreOffice install hint (Phase 2)

Output schema per slide:
  ## Slide N: {Title}
  [body text / bullet lists]
  [tables as GFM pipe tables]
  > Speaker notes (as blockquote)

Key design decisions:
- Title placeholder (idx 0) text goes into the Section heading, not the body
- Bullet paragraphs (level > 0, or buChar/buAutoNum XML markers) → IR List
- Tables in slide shapes → IR Table; first row treated as header
- Speaker notes appended as BlockQuote at end of each section
- Security: 50 MB input size limit; 500 MB zip bomb limit (.pptx is a ZIP)
- Metadata: all core properties from prs.core_properties
- Word count: sum of all text across all slides and notes

Install:
    pip install distill-core          # includes python-pptx
    pip install distill-core[legacy]  # adds .ppt support via LibreOffice (Phase 2)
"""

from __future__ import annotations

import io
import re
import zipfile
from pathlib import Path
from typing import Optional, Union

from distill.ir import (
    BlockQuote, Document, DocumentMetadata, Image, ImageType,
    List, ListItem, Paragraph, Section, Table, TableCell, TableRow, TextRun,
)
from distill.parsers.base import ParseError, ParseOptions, Parser, extract_image
from distill.registry import registry


# Placeholder indices for footer region shapes — suppress from output
_FOOTER_PLACEHOLDER_INDICES = frozenset({11, 12, 13})

# ── Security constants ────────────────────────────────────────────────────────

_MAX_FILE_BYTES  = 50  * 1024 * 1024   # 50 MB
_MAX_UNZIP_BYTES = 500 * 1024 * 1024   # 500 MB


def _check_input_size(source: Union[str, Path, bytes], max_bytes: int) -> None:
    """Raise ParseError if the source exceeds max_bytes."""
    mb = max_bytes // (1024 * 1024)
    if isinstance(source, (str, Path)):
        size = Path(source).stat().st_size
        if size > max_bytes:
            raise ParseError(
                f"Input file exceeds the {mb} MB size limit "
                f"({size / (1024*1024):.1f} MB). "
                f"Increase the limit via options.extra['max_file_size']."
            )
    elif isinstance(source, bytes):
        if len(source) > max_bytes:
            raise ParseError(
                f"Input file exceeds the {mb} MB size limit "
                f"({len(source) / (1024*1024):.1f} MB). "
                f"Increase the limit via options.extra['max_file_size']."
            )


def _check_zip_bomb(source: Union[str, Path, bytes], max_unzip_bytes: int) -> None:
    """Raise ParseError if the ZIP uncompressed size exceeds max_unzip_bytes."""
    mb = max_unzip_bytes // (1024 * 1024)
    try:
        data = source if isinstance(source, bytes) else Path(source).read_bytes()
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            total = sum(info.file_size for info in zf.infolist())
        if total > max_unzip_bytes:
            raise ParseError(
                f"PPTX archive uncompressed size ({total // (1024*1024)} MB) "
                f"exceeds the {mb} MB safety limit. "
                f"Increase via options.extra['max_unzip_size']."
            )
    except ParseError:
        raise
    except zipfile.BadZipFile:
        raise ParseError("File is not a valid PPTX (ZIP) archive.")
    except Exception as e:
        raise ParseError(f"Could not inspect archive: {e}") from e


# ── Metadata helpers ──────────────────────────────────────────────────────────

def _extract_metadata(prs, path: Optional[Path], slide_count: int) -> DocumentMetadata:
    """
    Pull core properties from prs.core_properties.
    python-pptx uses the same OOXML core-property model as python-docx.
    """
    cp = prs.core_properties

    def _safe(attr: str) -> Optional[str]:
        try:
            v = getattr(cp, attr, None)
            s = str(v).strip() if v is not None else None
            return s if s else None
        except Exception:
            return None

    def _iso(attr: str) -> Optional[str]:
        try:
            v = getattr(cp, attr, None)
            return v.isoformat() if v is not None else None
        except Exception:
            return None

    kw_raw   = _safe("keywords") or ""
    keywords = [k.strip() for k in re.split(r"[,;]", kw_raw) if k.strip()] if kw_raw else []

    # description: python-pptx uses .description (maps to dc:description)
    description = _safe("description") or _safe("comments") or None

    return DocumentMetadata(
        title         = _safe("title"),
        author        = _safe("author"),
        subject       = _safe("subject"),
        description   = description,
        keywords      = keywords,
        created_at    = _iso("created"),
        modified_at   = _iso("modified"),
        slide_count   = slide_count,
        source_format = "pptx",
        source_path   = str(path) if path else None,
    )


def _compute_word_count(prs) -> int:
    """Sum words across all slide text frames and speaker notes."""
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                total += len(shape.text_frame.text.split())
        if slide.has_notes_slide:
            total += len(slide.notes_slide.notes_text_frame.text.split())
    return total


def _get_first_font_size_pt(text_frame):
    """Return font size in points from the first run of the first non-empty
    paragraph in the text frame. Returns None if no explicit size is set."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                return run.font.size / 12700  # EMU to points
    return None


# ── Parser ────────────────────────────────────────────────────────────────────

@registry.register
class PptxParser(Parser):
    """Parses .pptx files using python-pptx."""

    extensions = [".pptx"]
    mime_types = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]
    requires          = ["pptx"]    # python-pptx (import name: pptx)
    optional_requires = []

    def parse(
        self,
        source: Union[str, Path, bytes],
        options: Optional[ParseOptions] = None,
    ) -> Document:
        options = options or ParseOptions()

        if isinstance(source, (str, Path)):
            path = Path(source)
        else:
            path = None

        # Security checks
        max_file  = options.extra.get("max_file_size",  _MAX_FILE_BYTES)
        max_unzip = options.extra.get("max_unzip_size", _MAX_UNZIP_BYTES)
        _check_input_size(source, max_file)
        _check_zip_bomb(source, max_unzip)

        try:
            from pptx import Presentation
        except ImportError as e:
            raise ParseError(f"Missing dependency: {e}") from e

        try:
            prs = Presentation(io.BytesIO(source) if isinstance(source, bytes) else str(source))
        except Exception as e:
            raise ParseError(f"python-pptx failed to open presentation: {e}") from e

        slide_count = len(prs.slides)
        metadata    = _extract_metadata(prs, path, slide_count)
        metadata.word_count = _compute_word_count(prs) or None
        document    = Document(metadata=metadata)

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_num, slide in enumerate(prs.slides, start=1):
            section = self._parse_slide(slide, slide_num, options, document, slide_width, slide_height)
            if section.blocks:
                document.sections.append(section)
            else:
                # Only suppress if heading is the generic "Slide N" fallback
                heading_text = section.heading[0].text if section.heading else ""
                if not re.match(r"^Slide \d+$", heading_text):
                    document.sections.append(section)

        return document

    # ── Slide parsing ─────────────────────────────────────────────────────────

    def _parse_slide(
        self,
        slide,
        slide_num: int,
        options: ParseOptions,
        doc: Document,
        slide_width: int = 0,
        slide_height: int = 0,
    ) -> Section:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from distill.parsers.base import classify_image

        title_text, title_shape_used = self._get_slide_title(slide, slide_height)
        heading = title_text if title_text else f"Slide {slide_num}"

        section = Section(
            heading=[TextRun(text=heading)],
            level=2,
        )

        title_placeholder = slide.shapes.title
        # Collect XML elements to skip — python-pptx creates new wrapper
        # objects on each access, so identity checks on shape objects fail.
        # Comparing the underlying lxml _element is stable.
        _skip_elements = set()
        if title_placeholder is not None:
            _skip_elements.add(id(title_placeholder._element))
        if title_shape_used is not None:
            _skip_elements.add(id(title_shape_used._element))

        source_stem = (
            Path(doc.metadata.source_path).stem
            if doc.metadata.source_path
            else "file"
        )
        image_index = 0

        for shape in slide.shapes:
            # Skip shapes whose text is already in the heading
            if id(shape._element) in _skip_elements:
                # Still parse remaining paragraphs if the shape has more
                # content beyond the title line (e.g. bullet lists)
                if shape.has_text_frame and len(shape.text_frame.paragraphs) > 1:
                    blocks = self._parse_text_frame(
                        shape.text_frame, skip_first=True,
                    )
                    section.blocks.extend(blocks)
                continue

            # Skip footer, date, and slide-number placeholders
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.idx in _FOOTER_PLACEHOLDER_INDICES:
                        continue
                except Exception:
                    pass  # placeholder_format may raise on malformed shapes

            if shape.has_text_frame:
                blocks = self._parse_text_frame(shape.text_frame)
                section.blocks.extend(blocks)

            elif shape.has_table:
                table = self._parse_table(shape.table, options)
                section.blocks.append(table)

            elif hasattr(shape, "shape_type"):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if options.images == "suppress":
                        continue

                    # Classify before blob read to avoid unnecessary ZIP reads
                    image_type = classify_image(
                        mode="pptx",
                        shape_w=shape.width,
                        shape_h=shape.height,
                        slide_w=slide_width,
                        slide_h=slide_height,
                        name=getattr(shape, "name", "") or "",
                    )

                    if image_type == ImageType.DECORATIVE:
                        doc.metadata.decorative_images_filtered += 1
                        continue

                    image_node = Image(
                        image_type=image_type,
                        alt_text=(
                            self._get_shape_descr(shape)
                            or (getattr(shape, "name", None) or "").strip()
                            or None
                        ),
                    )

                    # Extract image bytes from shape
                    try:
                        image_bytes = shape.image.blob
                        image_ext = getattr(shape.image, "ext", "png")
                    except Exception:
                        image_bytes = None
                        image_ext = "png"

                    if image_bytes and options.images in ("extract", "caption") and options.image_dir:
                        filename = f"{source_stem}_{slide_num}_{image_index}"
                        path = extract_image(
                            image_bytes=image_bytes,
                            ext=image_ext,
                            image_dir=Path(options.image_dir),
                            filename=filename,
                            collector=getattr(options, "collector", None),
                        )
                        if path:
                            image_node.path = path

                    # Vision captioning hookup
                    if options.images == "caption" and options.vision_provider and image_bytes:
                        try:
                            from distill.parsers._vision import caption_image
                            caption = caption_image(
                                image_bytes=image_bytes,
                                provider=options.vision_provider,
                                model=options.extra.get("vision_model") or None,
                                api_key=options.vision_api_key,
                                base_url=options.vision_base_url,
                            )
                            if caption:
                                image_node.caption = caption
                        except Exception as e:
                            if getattr(options, "collector", None) is not None:
                                from distill.warnings import ConversionWarning, WarningType
                                options.collector.add(ConversionWarning(
                                    type=WarningType.vision_caption_failed,
                                    message=f"Vision captioning failed for slide {slide_num} image {image_index}: {e}",
                                ))

                    section.blocks.append(image_node)
                    image_index += 1

        # Speaker notes appended as a block quote
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                bq = BlockQuote(content=[
                    Paragraph(runs=[TextRun(text=notes_text)])
                ])
                section.blocks.append(bq)

        return section

    def _get_slide_title(self, slide, slide_height: int = 0) -> tuple[Optional[str], Optional[object]]:
        # Primary: standard PowerPoint title placeholder
        try:
            title = slide.shapes.title
            if title and title.text and title.text.strip():
                return title.text.strip(), title
        except Exception:
            pass

        # Heuristic fallback: position + font size
        # Prefer shapes in the top 15% of the slide with font size >= 20pt
        if slide_height > 0:
            candidates = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text or len(text) > 120:
                    continue
                if shape.top is None:
                    continue
                # Skip bullet text frames
                if self._is_bullet_text_frame(shape):
                    continue
                top_ratio = shape.top / slide_height
                if top_ratio >= 0.15:
                    continue
                font_pt = _get_first_font_size_pt(shape.text_frame)
                if font_pt is None or font_pt < 20:
                    continue
                candidates.append((shape, font_pt))

            if candidates:
                # Largest font first, then highest on slide (smallest top)
                candidates.sort(key=lambda x: (-x[1], x[0].top))
                best_shape = candidates[0][0]
                first_line = best_shape.text_frame.text.strip().splitlines()[0].strip()
                return first_line, best_shape

        # Original z-order-first fallback (skip bullet text frames)
        try:
            from pptx.oxml.ns import qn as _qn
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if self._is_bullet_text_frame(shape):
                    continue
                first_line = text.splitlines()[0].strip()
                if first_line and len(first_line) <= 80:
                    return first_line, shape
        except Exception:
            pass

        return None, None

    def _is_bullet_text_frame(self, shape) -> bool:
        """Check if a shape's first paragraph is a bullet (used to skip in title extraction)."""
        try:
            from pptx.oxml.ns import qn as _qn
            first_para = shape.text_frame.paragraphs[0]
            _pPr = first_para._p.find(_qn("a:pPr"))
            if _pPr is not None:
                if (_pPr.find(_qn("a:buChar")) is not None
                        or _pPr.find(_qn("a:buAutoNum")) is not None):
                    if _pPr.find(_qn("a:buNone")) is None:
                        return True
        except Exception:
            pass
        return False

    def _parse_text_frame(self, text_frame, skip_first: bool = False) -> list:
        blocks: list = []
        # Accumulate (level, runs, ordered) tuples for consecutive bullet paragraphs
        bullet_acc: list[tuple[int, list, bool]] = []
        first_skipped = False

        for para in text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Skip the first non-empty paragraph (already used as slide title)
            if skip_first and not first_skipped:
                first_skipped = True
                continue

            runs = [
                TextRun(
                    text   = run.text,
                    bold   = bool(run.font.bold),
                    italic = bool(run.font.italic),
                )
                for run in para.runs
                if run.text
            ]
            if not runs:
                runs = [TextRun(text=text)]

            is_bullet, is_ordered = self._detect_bullet(para)

            if is_bullet:
                bullet_acc.append((para.level, runs, is_ordered))
            else:
                if bullet_acc:
                    blocks.append(self._build_list_from_flat(bullet_acc))
                    bullet_acc = []
                blocks.append(Paragraph(runs=runs))

        if bullet_acc:
            blocks.append(self._build_list_from_flat(bullet_acc))

        return blocks

    def _build_list_from_flat(
        self, items: list[tuple[int, list, bool]]
    ) -> List:
        """Convert flat (level, runs, ordered) tuples into a nested List IR."""
        # Determine ordered: if any item uses buAutoNum use ordered,
        # but if mixed, default to unordered
        all_ordered = all(o for _, _, o in items)
        any_ordered = any(o for _, _, o in items)
        ordered = all_ordered  # mixed = unordered (safe default)

        root = List(items=[], ordered=ordered)
        # Stack of (List node, level) for building nesting
        stack: list[tuple[List, int]] = [(root, 0)]

        for level, runs, _ in items:
            # Find or create the correct nesting depth
            while len(stack) > 1 and stack[-1][1] > level:
                stack.pop()

            parent_list = stack[-1][0]
            item = ListItem(content=runs)
            parent_list.items.append(item)

            # Prepare for potential children at deeper levels
            child_list = List(items=[], ordered=ordered)
            item.children.append(child_list)
            stack.append((child_list, level + 1))

        # Clean up empty child lists
        self._prune_empty_lists(root)
        return root

    def _prune_empty_lists(self, lst: List) -> None:
        """Remove empty child List nodes added speculatively."""
        for item in lst.items:
            item.children = [c for c in item.children if c.items]
            for child in item.children:
                self._prune_empty_lists(child)

    def _parse_table(self, tbl, options: ParseOptions) -> Table:
        rows: list[TableRow] = []
        for i, row in enumerate(tbl.rows):
            if options.max_table_rows > 0 and i >= options.max_table_rows:
                break
            cells = [
                TableCell(
                    content=[TextRun(text=cell.text.strip())],
                    is_header=(i == 0),
                )
                for cell in row.cells
            ]
            rows.append(TableRow(cells=cells))
        return Table(rows=rows)

    def _detect_bullet(self, para) -> tuple[bool, bool]:
        """
        Detect whether a paragraph is a bullet and whether it is ordered.

        Returns (is_bullet, is_ordered).  Checks <a:buChar> and <a:buAutoNum>
        inside <a:pPr>, and also treats para.level > 0 as a bullet (indented
        paragraphs in content placeholders).  <a:buNone> suppresses bullets.
        """
        try:
            from pptx.oxml.ns import qn
            pPr = para._p.find(qn("a:pPr"))
            if pPr is not None:
                has_buNone = pPr.find(qn("a:buNone")) is not None
                if has_buNone:
                    return (False, False)
                has_buChar = pPr.find(qn("a:buChar")) is not None
                has_buAutoNum = pPr.find(qn("a:buAutoNum")) is not None
                if has_buChar or has_buAutoNum:
                    is_ordered = has_buAutoNum and not has_buChar
                    return (True, is_ordered)
            # Indented paragraphs (level > 0) are also treated as bullets
            if para.level > 0:
                return (True, False)
            return (False, False)
        except Exception:
            return (False, False)

    def _get_shape_descr(self, shape) -> str:
        """Read the author-provided alt text (descr attribute) from a shape's XML.
        Returns stripped string or empty string if not available."""
        try:
            cNvPr = shape._element.nvPicPr.cNvPr
            descr = cNvPr.get("descr") or ""
            return descr.strip()
        except Exception:
            return ""


# ── Legacy parser ─────────────────────────────────────────────────────────────

@registry.register
class PptLegacyParser(Parser):
    """
    Converts legacy .ppt binary presentations to .pptx via LibreOffice headless,
    then delegates to PptxParser for the actual content extraction.

    Requires LibreOffice to be installed and on PATH (or DISTILL_LIBREOFFICE
    environment variable set to the full binary path).
    """

    extensions            = [".ppt"]
    mime_types            = ["application/vnd.ms-powerpoint"]
    requires              = ["pptx"]
    requires_libreoffice  = True

    def parse(
        self,
        source: Union[str, Path, bytes],
        options: Optional[ParseOptions] = None,
    ) -> Document:
        import shutil
        from distill.parsers._libreoffice import convert_via_libreoffice

        options = options or ParseOptions()
        timeout = options.extra.get("libreoffice_timeout", 60)

        output_path = convert_via_libreoffice(source, "pptx", timeout=timeout)
        try:
            doc = PptxParser().parse(output_path, options)
            # Preserve the original source path/format in metadata
            doc.metadata.source_format = "ppt"
            if not isinstance(source, bytes):
                doc.metadata.source_path = str(source)
            return doc
        finally:
            shutil.rmtree(output_path.parent, ignore_errors=True)
