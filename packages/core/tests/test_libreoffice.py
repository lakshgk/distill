"""
Tests for distill.parsers._libreoffice and legacy format parsers.

All LibreOffice subprocess calls are mocked — these tests verify the
integration logic without requiring LibreOffice to be installed.
"""

from __future__ import annotations

import io
import os
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Optional
from unittest.mock import MagicMock, patch

import pytest

from distill.parsers._libreoffice import (
    convert_via_libreoffice,
    find_libreoffice,
    is_libreoffice_available,
)
from distill.parsers.base import ParseError


# ── Helpers ───────────────────────────────────────────────────────────────────

def _make_xlsx_bytes() -> bytes:
    """Minimal valid .xlsx (ZIP) for delegation tests."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value"])
    ws.append(["Alpha", "1"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx_bytes() -> bytes:
    """Minimal valid .pptx (ZIP) for delegation tests."""
    import pptx
    from pptx.util import Inches
    prs = pptx.Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test Slide"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_docx_bytes() -> bytes:
    """Minimal valid .docx (ZIP) for delegation tests."""
    import docx
    doc = docx.Document()
    doc.add_paragraph("Hello world")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _write_converted_file(tmp_dir: Path, ext: str, content: bytes) -> None:
    """Write a fake converted file into tmp_dir, as LibreOffice would."""
    (tmp_dir / f"input.{ext}").write_bytes(content)


# ── find_libreoffice ─────────────────────────────────────────────────────────

class TestFindLibreoffice:
    def test_returns_none_when_nothing_available(self):
        with patch.dict(os.environ, {}, clear=True):
            # Remove DISTILL_LIBREOFFICE if set
            os.environ.pop("DISTILL_LIBREOFFICE", None)
            with patch("shutil.which", return_value=None), \
                 patch("pathlib.Path.is_file", return_value=False):
                result = find_libreoffice()
                assert result is None

    def test_env_override_takes_priority(self, tmp_path):
        fake_bin = tmp_path / "soffice"
        fake_bin.touch()
        with patch.dict(os.environ, {"DISTILL_LIBREOFFICE": str(fake_bin)}):
            result = find_libreoffice()
            assert result == str(fake_bin)

    def test_env_override_missing_file_falls_through(self, tmp_path):
        nonexistent = str(tmp_path / "does_not_exist")
        with patch.dict(os.environ, {"DISTILL_LIBREOFFICE": nonexistent}):
            with patch("shutil.which", return_value=None), \
                 patch("pathlib.Path.is_file", return_value=False):
                result = find_libreoffice()
                assert result is None

    def test_finds_via_which(self):
        with patch.dict(os.environ, {}, clear=True):
            os.environ.pop("DISTILL_LIBREOFFICE", None)
            with patch("shutil.which", side_effect=lambda x: "/usr/bin/libreoffice" if x == "libreoffice" else None):
                result = find_libreoffice()
                assert result == "/usr/bin/libreoffice"

    def test_finds_via_absolute_candidate_path(self, tmp_path):
        fake_lo = tmp_path / "soffice"
        fake_lo.touch()
        with patch.dict(os.environ, {}, clear=True):
            os.environ.pop("DISTILL_LIBREOFFICE", None)
            with patch("shutil.which", return_value=None):
                from distill.parsers import _libreoffice as lo_mod
                original = lo_mod._CANDIDATE_PATHS
                lo_mod._CANDIDATE_PATHS = [str(fake_lo)]
                try:
                    result = find_libreoffice()
                    assert result == str(fake_lo)
                finally:
                    lo_mod._CANDIDATE_PATHS = original


class TestIsLibreofficeAvailable:
    def test_returns_true_when_found(self):
        with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/usr/bin/soffice"):
            assert is_libreoffice_available() is True

    def test_returns_false_when_not_found(self):
        with patch("distill.parsers._libreoffice.find_libreoffice", return_value=None):
            assert is_libreoffice_available() is False


# ── convert_via_libreoffice ───────────────────────────────────────────────────

class TestConvertViaLibreoffice:
    def _mock_successful_run(self, tmp_dir: Path, ext: str, content: bytes):
        """
        Return a mock for subprocess.run that also writes the expected
        output file into tmp_dir so the path-existence check passes.
        """
        def side_effect(cmd, **kwargs):
            # Write the fake converted file
            stem = Path(cmd[-1]).stem
            (tmp_dir / f"{stem}.{ext}").write_bytes(content)
            return MagicMock(returncode=0, stdout="", stderr="")
        return side_effect

    def test_successful_conversion_from_path(self, tmp_path):
        src = tmp_path / "doc.doc"
        src.write_bytes(b"fake doc content")
        xlsx_bytes = _make_xlsx_bytes()

        # Create real temp dir BEFORE patching so mkdtemp isn't intercepted here
        real_tmp = tempfile.mkdtemp(prefix="test_lo_")
        try:
            with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/usr/bin/soffice"), \
                 patch("distill.parsers._libreoffice.tempfile.mkdtemp", return_value=real_tmp), \
                 patch("subprocess.run") as mock_run:
                mock_run.side_effect = self._mock_successful_run(
                    Path(real_tmp), "xlsx", xlsx_bytes
                )
                result = convert_via_libreoffice(src, "xlsx")
                assert result.exists()
                assert result.suffix == ".xlsx"
        finally:
            shutil.rmtree(real_tmp, ignore_errors=True)

    def test_successful_conversion_from_bytes(self):
        src_bytes = b"fake doc binary content"
        xlsx_bytes = _make_xlsx_bytes()

        real_tmp = tempfile.mkdtemp(prefix="test_lo_")
        try:
            with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/usr/bin/soffice"), \
                 patch("distill.parsers._libreoffice.tempfile.mkdtemp", return_value=real_tmp), \
                 patch("subprocess.run") as mock_run:
                mock_run.side_effect = self._mock_successful_run(
                    Path(real_tmp), "xlsx", xlsx_bytes
                )
                result = convert_via_libreoffice(src_bytes, "xlsx")
                assert result.exists()
                assert result.suffix == ".xlsx"
        finally:
            shutil.rmtree(real_tmp, ignore_errors=True)

    def test_raises_when_libreoffice_not_found(self, tmp_path):
        src = tmp_path / "file.doc"
        src.write_bytes(b"x")
        with patch("distill.parsers._libreoffice.find_libreoffice", return_value=None):
            with pytest.raises(ParseError, match="LibreOffice is not installed"):
                convert_via_libreoffice(src, "docx")

    def test_raises_on_timeout(self, tmp_path):
        src = tmp_path / "file.doc"
        src.write_bytes(b"x")
        with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/usr/bin/soffice"), \
             patch("subprocess.run", side_effect=subprocess.TimeoutExpired(cmd="soffice", timeout=5)):
            with pytest.raises(ParseError, match="timed out"):
                convert_via_libreoffice(src, "docx", timeout=5)

    def test_raises_on_nonzero_exit(self, tmp_path):
        src = tmp_path / "file.doc"
        src.write_bytes(b"x")
        real_tmp = tempfile.mkdtemp(prefix="test_lo_")
        try:
            with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/usr/bin/soffice"), \
                 patch("distill.parsers._libreoffice.tempfile.mkdtemp", return_value=real_tmp), \
                 patch("subprocess.run", return_value=MagicMock(
                     returncode=1,
                     stdout="",
                     stderr="conversion failed: bad format",
                 )):
                with pytest.raises(ParseError, match="conversion failed"):
                    convert_via_libreoffice(src, "docx")
        finally:
            shutil.rmtree(real_tmp, ignore_errors=True)

    def test_raises_when_output_file_missing(self, tmp_path):
        src = tmp_path / "file.doc"
        src.write_bytes(b"x")
        real_tmp = tempfile.mkdtemp(prefix="test_lo_")
        try:
            with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/usr/bin/soffice"), \
                 patch("distill.parsers._libreoffice.tempfile.mkdtemp", return_value=real_tmp), \
                 patch("subprocess.run", return_value=MagicMock(
                     returncode=0, stdout="", stderr=""
                 )):
                # Run succeeds but writes no output file — real_tmp is empty
                with pytest.raises(ParseError, match="produced no .docx file"):
                    convert_via_libreoffice(src, "docx")
        finally:
            shutil.rmtree(real_tmp, ignore_errors=True)

    def test_raises_on_binary_not_found(self, tmp_path):
        src = tmp_path / "file.doc"
        src.write_bytes(b"x")
        real_tmp = tempfile.mkdtemp(prefix="test_lo_")
        try:
            with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/fake/path/soffice"), \
                 patch("distill.parsers._libreoffice.tempfile.mkdtemp", return_value=real_tmp), \
                 patch("subprocess.run", side_effect=FileNotFoundError()):
                with pytest.raises(ParseError, match="not found at"):
                    convert_via_libreoffice(src, "docx")
        finally:
            shutil.rmtree(real_tmp, ignore_errors=True)

    def test_cleans_up_tmpdir_on_error(self, tmp_path):
        src = tmp_path / "file.doc"
        src.write_bytes(b"x")
        captured_tmp = []

        original_mkdtemp = tempfile.mkdtemp

        def capture_mkdtemp(**kwargs):
            d = original_mkdtemp(**kwargs)
            captured_tmp.append(d)
            return d

        with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/usr/bin/soffice"), \
             patch("distill.parsers._libreoffice.tempfile.mkdtemp", side_effect=capture_mkdtemp), \
             patch("subprocess.run", side_effect=subprocess.TimeoutExpired(cmd="soffice", timeout=5)):
            with pytest.raises(ParseError):
                convert_via_libreoffice(src, "docx")

        # Temp dir should have been cleaned up on error
        for d in captured_tmp:
            assert not Path(d).exists(), f"Temp dir {d} was not cleaned up"

    def test_command_includes_user_installation_flag(self, tmp_path):
        src = tmp_path / "file.doc"
        src.write_bytes(b"x")
        xlsx_bytes = _make_xlsx_bytes()

        real_tmp = tempfile.mkdtemp(prefix="test_lo_")
        try:
            captured_cmd = []

            def capture_run(cmd, **kwargs):
                captured_cmd.extend(cmd)
                stem = Path(cmd[-1]).stem
                (Path(real_tmp) / f"{stem}.xlsx").write_bytes(xlsx_bytes)
                return MagicMock(returncode=0, stdout="", stderr="")

            with patch("distill.parsers._libreoffice.find_libreoffice", return_value="/usr/bin/soffice"), \
                 patch("distill.parsers._libreoffice.tempfile.mkdtemp", return_value=real_tmp), \
                 patch("subprocess.run", side_effect=capture_run):
                convert_via_libreoffice(src, "xlsx")

            assert "--headless" in captured_cmd
            assert "--norestore" in captured_cmd
            assert any("-env:UserInstallation" in arg for arg in captured_cmd)
        finally:
            shutil.rmtree(real_tmp, ignore_errors=True)


# ── DocLegacyParser ──────────────────────────────────────────────────────────

class TestDocLegacyParser:
    def test_delegates_to_docx_parser(self, tmp_path):
        from distill.parsers.docx import DocLegacyParser

        src = tmp_path / "legacy.doc"
        src.write_bytes(b"fake legacy doc")
        docx_bytes = _make_docx_bytes()

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_doc_lo_"))
            out_file = out_dir / "legacy.docx"
            out_file.write_bytes(docx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            doc = DocLegacyParser().parse(src)

        assert doc.metadata.source_format == "doc"
        assert doc.metadata.source_path == str(src)

    def test_source_format_overridden_to_doc(self, tmp_path):
        from distill.parsers.docx import DocLegacyParser

        src = tmp_path / "legacy.doc"
        src.write_bytes(b"fake legacy doc")
        docx_bytes = _make_docx_bytes()

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_doc_lo_"))
            out_file = out_dir / "legacy.docx"
            out_file.write_bytes(docx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            doc = DocLegacyParser().parse(src)

        # Must be "doc" not "docx" (the docx parser would set "docx")
        assert doc.metadata.source_format == "doc"

    def test_cleans_up_tmpdir_after_success(self, tmp_path):
        from distill.parsers.docx import DocLegacyParser

        src = tmp_path / "legacy.doc"
        src.write_bytes(b"fake legacy doc")
        docx_bytes = _make_docx_bytes()
        created_dirs = []

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_doc_lo_"))
            created_dirs.append(out_dir)
            out_file = out_dir / "legacy.docx"
            out_file.write_bytes(docx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            DocLegacyParser().parse(src)

        for d in created_dirs:
            assert not d.exists(), f"Temp dir {d} was not cleaned up"

    def test_cleans_up_tmpdir_on_parse_error(self, tmp_path):
        from distill.parsers.docx import DocLegacyParser
        from distill.parsers.base import ParseError as PE

        src = tmp_path / "legacy.doc"
        src.write_bytes(b"fake legacy doc")
        created_dirs = []

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_doc_lo_"))
            created_dirs.append(out_dir)
            # Write corrupt (empty) docx — DocxParser will raise
            out_file = out_dir / "legacy.docx"
            out_file.write_bytes(b"not a real docx")
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            with pytest.raises(Exception):
                DocLegacyParser().parse(src)

        for d in created_dirs:
            assert not d.exists(), f"Temp dir {d} was not cleaned up on error"

    def test_libreoffice_timeout_option_forwarded(self, tmp_path):
        from distill.parsers.docx import DocLegacyParser
        from distill.parsers.base import ParseOptions

        src = tmp_path / "legacy.doc"
        src.write_bytes(b"fake")
        received_timeout = []

        def fake_convert(source, target_ext, timeout=60):
            received_timeout.append(timeout)
            raise ParseError("abort early")

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            with pytest.raises(ParseError):
                opts = ParseOptions()
                opts.extra["libreoffice_timeout"] = 120
                DocLegacyParser().parse(src, opts)

        assert received_timeout == [120]


# ── XlsLegacyParser ──────────────────────────────────────────────────────────

class TestXlsLegacyParser:
    def test_delegates_to_xlsx_parser(self, tmp_path):
        from distill.parsers.xlsx import XlsLegacyParser

        src = tmp_path / "legacy.xls"
        src.write_bytes(b"fake legacy xls")
        xlsx_bytes = _make_xlsx_bytes()

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_xls_lo_"))
            out_file = out_dir / "legacy.xlsx"
            out_file.write_bytes(xlsx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            doc = XlsLegacyParser().parse(src)

        assert doc.metadata.source_format == "xls"
        assert doc.metadata.source_path == str(src)

    def test_source_format_overridden_to_xls(self, tmp_path):
        from distill.parsers.xlsx import XlsLegacyParser

        src = tmp_path / "legacy.xls"
        src.write_bytes(b"fake legacy xls")
        xlsx_bytes = _make_xlsx_bytes()

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_xls_lo_"))
            out_file = out_dir / "legacy.xlsx"
            out_file.write_bytes(xlsx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            doc = XlsLegacyParser().parse(src)

        assert doc.metadata.source_format == "xls"

    def test_cleans_up_tmpdir_after_success(self, tmp_path):
        from distill.parsers.xlsx import XlsLegacyParser

        src = tmp_path / "legacy.xls"
        src.write_bytes(b"fake legacy xls")
        xlsx_bytes = _make_xlsx_bytes()
        created_dirs = []

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_xls_lo_"))
            created_dirs.append(out_dir)
            out_file = out_dir / "legacy.xlsx"
            out_file.write_bytes(xlsx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            XlsLegacyParser().parse(src)

        for d in created_dirs:
            assert not d.exists(), f"Temp dir {d} was not cleaned up"

    def test_libreoffice_timeout_option_forwarded(self, tmp_path):
        from distill.parsers.xlsx import XlsLegacyParser
        from distill.parsers.base import ParseOptions

        src = tmp_path / "legacy.xls"
        src.write_bytes(b"fake")
        received_timeout = []

        def fake_convert(source, target_ext, timeout=60):
            received_timeout.append(timeout)
            raise ParseError("abort early")

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            with pytest.raises(ParseError):
                opts = ParseOptions()
                opts.extra["libreoffice_timeout"] = 90
                XlsLegacyParser().parse(src, opts)

        assert received_timeout == [90]

    def test_produced_document_has_sections(self, tmp_path):
        from distill.parsers.xlsx import XlsLegacyParser

        src = tmp_path / "legacy.xls"
        src.write_bytes(b"fake legacy xls")
        xlsx_bytes = _make_xlsx_bytes()

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_xls_lo_"))
            out_file = out_dir / "legacy.xlsx"
            out_file.write_bytes(xlsx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            doc = XlsLegacyParser().parse(src)

        assert len(doc.sections) > 0


# ── PptLegacyParser ──────────────────────────────────────────────────────────

class TestPptLegacyParser:
    def test_delegates_to_pptx_parser(self, tmp_path):
        from distill.parsers.pptx import PptLegacyParser

        src = tmp_path / "legacy.ppt"
        src.write_bytes(b"fake legacy ppt")
        pptx_bytes = _make_pptx_bytes()

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_ppt_lo_"))
            out_file = out_dir / "legacy.pptx"
            out_file.write_bytes(pptx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            doc = PptLegacyParser().parse(src)

        assert doc.metadata.source_format == "ppt"
        assert doc.metadata.source_path == str(src)

    def test_source_format_overridden_to_ppt(self, tmp_path):
        from distill.parsers.pptx import PptLegacyParser

        src = tmp_path / "legacy.ppt"
        src.write_bytes(b"fake legacy ppt")
        pptx_bytes = _make_pptx_bytes()

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_ppt_lo_"))
            out_file = out_dir / "legacy.pptx"
            out_file.write_bytes(pptx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            doc = PptLegacyParser().parse(src)

        assert doc.metadata.source_format == "ppt"

    def test_cleans_up_tmpdir_after_success(self, tmp_path):
        from distill.parsers.pptx import PptLegacyParser

        src = tmp_path / "legacy.ppt"
        src.write_bytes(b"fake legacy ppt")
        pptx_bytes = _make_pptx_bytes()
        created_dirs = []

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_ppt_lo_"))
            created_dirs.append(out_dir)
            out_file = out_dir / "legacy.pptx"
            out_file.write_bytes(pptx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            PptLegacyParser().parse(src)

        for d in created_dirs:
            assert not d.exists(), f"Temp dir {d} was not cleaned up"

    def test_libreoffice_timeout_option_forwarded(self, tmp_path):
        from distill.parsers.pptx import PptLegacyParser
        from distill.parsers.base import ParseOptions

        src = tmp_path / "legacy.ppt"
        src.write_bytes(b"fake")
        received_timeout = []

        def fake_convert(source, target_ext, timeout=60):
            received_timeout.append(timeout)
            raise ParseError("abort early")

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            with pytest.raises(ParseError):
                opts = ParseOptions()
                opts.extra["libreoffice_timeout"] = 45
                PptLegacyParser().parse(src, opts)

        assert received_timeout == [45]

    def test_produced_document_has_sections(self, tmp_path):
        from distill.parsers.pptx import PptLegacyParser

        src = tmp_path / "legacy.ppt"
        src.write_bytes(b"fake legacy ppt")
        pptx_bytes = _make_pptx_bytes()

        def fake_convert(source, target_ext, timeout=60):
            out_dir = Path(tempfile.mkdtemp(prefix="test_ppt_lo_"))
            out_file = out_dir / "legacy.pptx"
            out_file.write_bytes(pptx_bytes)
            return out_file

        with patch("distill.parsers._libreoffice.convert_via_libreoffice", side_effect=fake_convert):
            doc = PptLegacyParser().parse(src)

        assert len(doc.sections) > 0
