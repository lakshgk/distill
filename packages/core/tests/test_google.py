"""
Tests for distill.parsers.google — Google Workspace parsers.

All Google API calls are mocked.  No network access is made and the
google-api-python-client / google-auth packages are not required to be
installed in order to run these tests.

Strategy
--------
The parsers use lazy imports inside their functions, so we can patch
``sys.modules`` to inject MagicMock stubs for all Google libraries before
calling any parser code.
"""

from __future__ import annotations

import io
import json
import os
import sys
import tempfile
from pathlib import Path
from types import ModuleType
from typing import Any
from unittest.mock import MagicMock, patch, PropertyMock

import pytest

from distill.parsers.base import ParseError, ParseOptions
from distill.parsers.base import UnsupportedFormatError
from distill.parsers.google import (
    GoogleDocsParser,
    GoogleSheetsParser,
    GoogleSlidesParser,
    _EXPORT_MAP,
    _REQUIRES,
    _build_credentials,
    _export_via_drive,
    _extract_file_id,
)


# ── Fake-module fixtures ───────────────────────────────────────────────────────

FAKE_FILE_ID = "1BxiMVs0XRA5nFMdKvBdBZjgmUUqptlbs74OgVE2upms"

DOCS_URL   = f"https://docs.google.com/document/d/{FAKE_FILE_ID}/edit"
SHEETS_URL = f"https://docs.google.com/spreadsheets/d/{FAKE_FILE_ID}/edit"
SLIDES_URL = f"https://docs.google.com/presentation/d/{FAKE_FILE_ID}/edit"


def _make_google_stubs():
    """
    Return a dict of ``sys.modules`` patches that provide minimal Google
    library stubs.  Only the interfaces actually used by google.py are mocked.
    """
    # google.oauth2.credentials.Credentials
    mock_credentials_cls = MagicMock()
    mock_credentials_obj = MagicMock()
    mock_credentials_obj.token = "fake-token"
    mock_credentials_cls.return_value = mock_credentials_obj

    # google.oauth2.service_account.Credentials
    mock_sa_cls = MagicMock()
    mock_sa_obj = MagicMock()
    mock_sa_obj.service_account_email = "svc@project.iam.gserviceaccount.com"
    mock_sa_cls.from_service_account_file.return_value = mock_sa_obj

    mock_google            = MagicMock()
    mock_google_oauth2     = MagicMock()
    mock_google_oauth2_creds   = MagicMock()
    mock_google_oauth2_sa      = MagicMock()

    mock_google_oauth2_creds.Credentials         = mock_credentials_cls
    mock_google_oauth2_sa.Credentials            = mock_sa_cls
    mock_google_oauth2.credentials               = mock_google_oauth2_creds
    mock_google_oauth2.service_account           = mock_google_oauth2_sa

    mock_google.oauth2 = mock_google_oauth2
    mock_google.auth   = MagicMock()

    mock_googleapiclient         = MagicMock()
    mock_googleapiclient_disco   = MagicMock()
    mock_googleapiclient_http    = MagicMock()

    return {
        "google":                       mock_google,
        "google.oauth2":                mock_google_oauth2,
        "google.oauth2.credentials":    mock_google_oauth2_creds,
        "google.oauth2.service_account": mock_google_oauth2_sa,
        "google.auth":                  mock_google.auth,
        "googleapiclient":              mock_googleapiclient,
        "googleapiclient.discovery":    mock_googleapiclient_disco,
        "googleapiclient.http":         mock_googleapiclient_http,
    }, {
        "credentials_cls":  mock_credentials_cls,
        "credentials_obj":  mock_credentials_obj,
        "sa_cls":           mock_sa_cls,
        "sa_obj":           mock_sa_obj,
        "discovery":        mock_googleapiclient_disco,
        "http":             mock_googleapiclient_http,
    }


def _make_drive_service(
    mime_type:  str,
    file_name:  str,
    export_bytes: bytes,
    stubs: dict,
):
    """
    Configure the *stubs* dict so that ``googleapiclient.discovery.build``
    returns a Drive service mock that:
      - returns *mime_type* and *file_name* from ``files().get()``
      - returns *export_bytes* from ``files().export_media()``
    """
    # Build a chain:  service.files().get(...).execute() → meta
    mock_service   = MagicMock()
    mock_files     = MagicMock()
    mock_get_req   = MagicMock()
    mock_export_req = MagicMock()

    meta = {"mimeType": mime_type, "name": file_name}
    mock_get_req.execute.return_value = meta
    mock_files.get.return_value       = mock_get_req

    mock_files.export_media.return_value = mock_export_req

    mock_service.files.return_value = mock_files
    stubs["discovery"].build.return_value = mock_service

    # MediaIoBaseDownload: on first next_chunk call write bytes and return done=True
    def make_downloader(buf, request):
        dl = MagicMock()
        def next_chunk():
            buf.write(export_bytes)
            return (None, True)
        dl.next_chunk.side_effect = next_chunk
        return dl

    stubs["http"].MediaIoBaseDownload.side_effect = make_downloader

    return mock_service


def _make_docx_bytes() -> bytes:
    """Minimal valid .docx."""
    import docx as _docx
    doc = _docx.Document()
    doc.add_paragraph("Hello from Google Docs")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes() -> bytes:
    """Minimal valid .xlsx."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score"])
    ws.append(["Alice", "95"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx_bytes() -> bytes:
    """Minimal valid .pptx."""
    import pptx as _pptx
    prs = _pptx.Presentation()
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test Slide"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── _extract_file_id ──────────────────────────────────────────────────────────

class TestExtractFileId:
    def test_docs_edit_url(self):
        assert _extract_file_id(DOCS_URL) == FAKE_FILE_ID

    def test_sheets_url_with_gid(self):
        url = f"https://docs.google.com/spreadsheets/d/{FAKE_FILE_ID}/edit#gid=0"
        assert _extract_file_id(url) == FAKE_FILE_ID

    def test_slides_url(self):
        assert _extract_file_id(SLIDES_URL) == FAKE_FILE_ID

    def test_drive_view_url(self):
        url = f"https://drive.google.com/file/d/{FAKE_FILE_ID}/view?usp=sharing"
        assert _extract_file_id(url) == FAKE_FILE_ID

    def test_bare_file_id(self):
        assert _extract_file_id(FAKE_FILE_ID) == FAKE_FILE_ID

    def test_gdoc_shortcut_path(self, tmp_path):
        p = tmp_path / f"{FAKE_FILE_ID}.gdoc"
        p.touch()
        assert _extract_file_id(str(p)) == FAKE_FILE_ID

    def test_gsheet_shortcut_path(self, tmp_path):
        p = tmp_path / f"{FAKE_FILE_ID}.gsheet"
        p.touch()
        assert _extract_file_id(str(p)) == FAKE_FILE_ID

    def test_gslides_shortcut_path(self, tmp_path):
        p = tmp_path / f"{FAKE_FILE_ID}.gslides"
        p.touch()
        assert _extract_file_id(str(p)) == FAKE_FILE_ID

    def test_returns_none_for_garbage(self):
        assert _extract_file_id("not-a-url") is None

    def test_returns_none_for_non_google_url(self):
        assert _extract_file_id("https://example.com/doc/123") is None

    def test_returns_none_for_empty_string(self):
        assert _extract_file_id("") is None

    def test_returns_none_for_short_id(self):
        # IDs must be 25+ chars
        assert _extract_file_id("shortid") is None

    def test_strips_whitespace(self):
        assert _extract_file_id(f"  {DOCS_URL}  ") == FAKE_FILE_ID


# ── _build_credentials ────────────────────────────────────────────────────────

class TestBuildCredentials:
    def test_access_token_builds_credentials(self):
        stubs, refs = _make_google_stubs()
        opts = ParseOptions()
        opts.extra["access_token"] = "ya29.token"

        with patch.dict(sys.modules, stubs):
            creds = _build_credentials(opts)

        refs["credentials_cls"].assert_called_once_with(token="ya29.token")

    def test_service_account_path_string(self, tmp_path):
        stubs, refs = _make_google_stubs()
        sa_file = tmp_path / "sa.json"
        sa_file.write_text(json.dumps({"type": "service_account"}))

        opts = ParseOptions()
        opts.extra["google_credentials"] = str(sa_file)

        with patch.dict(sys.modules, stubs):
            creds = _build_credentials(opts)

        refs["sa_cls"].from_service_account_file.assert_called_once_with(
            str(sa_file), scopes=["https://www.googleapis.com/auth/drive.readonly"]
        )

    def test_credentials_object_returned_as_is(self):
        stubs, refs = _make_google_stubs()
        opts = ParseOptions()
        # Simulate a pre-built Credentials object (has a 'token' attribute)
        fake_creds = MagicMock()
        fake_creds.token = "prebuilt-token"
        opts.extra["google_credentials"] = fake_creds

        with patch.dict(sys.modules, stubs):
            result = _build_credentials(opts)

        assert result is fake_creds

    def test_service_account_via_env_var(self, tmp_path, monkeypatch):
        stubs, refs = _make_google_stubs()
        sa_file = tmp_path / "sa.json"
        sa_file.write_text(json.dumps({"type": "service_account"}))
        monkeypatch.setenv("DISTILL_GOOGLE_CREDENTIALS", str(sa_file))

        opts = ParseOptions()

        with patch.dict(sys.modules, stubs):
            _build_credentials(opts)

        refs["sa_cls"].from_service_account_file.assert_called_once_with(
            str(sa_file), scopes=["https://www.googleapis.com/auth/drive.readonly"]
        )

    def test_raises_when_no_credentials(self, monkeypatch):
        stubs, refs = _make_google_stubs()
        monkeypatch.delenv("DISTILL_GOOGLE_CREDENTIALS", raising=False)
        opts = ParseOptions()

        with patch.dict(sys.modules, stubs):
            with pytest.raises(ParseError, match="requires credentials"):
                _build_credentials(opts)

    def test_raises_on_bad_service_account_path(self, tmp_path):
        stubs, refs = _make_google_stubs()
        refs["sa_cls"].from_service_account_file.side_effect = ValueError("bad key")

        opts = ParseOptions()
        opts.extra["google_credentials"] = str(tmp_path / "missing.json")

        with patch.dict(sys.modules, stubs):
            with pytest.raises(ParseError, match="Could not load"):
                _build_credentials(opts)

    def test_raises_when_google_auth_not_installed(self):
        opts = ParseOptions()
        opts.extra["access_token"] = "ya29.token"

        # Remove google.oauth2 from sys.modules entirely so import raises ImportError
        with patch.dict(sys.modules, {"google.oauth2": None, "google.oauth2.credentials": None,
                                      "google.oauth2.service_account": None}):
            with pytest.raises(ParseError, match="not available"):
                _build_credentials(opts)


# ── _export_via_drive ─────────────────────────────────────────────────────────

class TestExportViaDrive:
    def test_successful_export(self):
        stubs, refs = _make_google_stubs()
        docx_bytes  = _make_docx_bytes()
        fake_creds  = MagicMock()

        _make_drive_service(
            mime_type    = "application/vnd.google-apps.document",
            file_name    = "My Doc",
            export_bytes = docx_bytes,
            stubs        = refs,
        )

        with patch.dict(sys.modules, stubs):
            content, name = _export_via_drive(
                FAKE_FILE_ID,
                "application/vnd.google-apps.document",
                fake_creds,
                DOCS_URL,
            )

        assert content == docx_bytes
        assert name == "My Doc"

    def test_raises_on_wrong_mime_type(self):
        stubs, refs = _make_google_stubs()
        fake_creds  = MagicMock()

        _make_drive_service(
            mime_type    = "application/vnd.google-apps.spreadsheet",  # wrong!
            file_name    = "A Sheet",
            export_bytes = b"",
            stubs        = refs,
        )

        with patch.dict(sys.modules, stubs):
            with pytest.raises(UnsupportedFormatError, match="Expected a Google"):
                _export_via_drive(
                    FAKE_FILE_ID,
                    "application/vnd.google-apps.document",   # expected Docs
                    fake_creds,
                    DOCS_URL,
                )

    def test_raises_on_permission_denied(self):
        stubs, refs = _make_google_stubs()
        fake_creds  = MagicMock()

        # Simulate a 403 response
        http_error = Exception("HttpError 403: The caller does not have permission")
        refs["discovery"].build.return_value.files.return_value.get.return_value.execute.side_effect = http_error

        with patch.dict(sys.modules, stubs):
            with pytest.raises(ParseError, match="Permission denied"):
                _export_via_drive(
                    FAKE_FILE_ID,
                    "application/vnd.google-apps.document",
                    fake_creds,
                    DOCS_URL,
                )

    def test_raises_on_not_found(self):
        stubs, refs = _make_google_stubs()
        fake_creds  = MagicMock()

        http_error = Exception("HttpError 404: File not found")
        refs["discovery"].build.return_value.files.return_value.get.return_value.execute.side_effect = http_error

        with patch.dict(sys.modules, stubs):
            with pytest.raises(ParseError, match="not found"):
                _export_via_drive(
                    FAKE_FILE_ID,
                    "application/vnd.google-apps.document",
                    fake_creds,
                    DOCS_URL,
                )

    def test_raises_on_generic_api_error(self):
        stubs, refs = _make_google_stubs()
        fake_creds  = MagicMock()

        refs["discovery"].build.return_value.files.return_value.get.return_value.execute.side_effect = RuntimeError("service unavailable")

        with patch.dict(sys.modules, stubs):
            with pytest.raises(ParseError, match="Google Drive API error"):
                _export_via_drive(
                    FAKE_FILE_ID,
                    "application/vnd.google-apps.document",
                    fake_creds,
                    DOCS_URL,
                )

    def test_raises_when_googleapiclient_not_installed(self):
        fake_creds = MagicMock()
        with patch.dict(sys.modules, {"googleapiclient": None,
                                      "googleapiclient.discovery": None,
                                      "googleapiclient.http": None}):
            with pytest.raises(ParseError, match="not available"):
                _export_via_drive(
                    FAKE_FILE_ID,
                    "application/vnd.google-apps.document",
                    fake_creds,
                    DOCS_URL,
                )


# ── GoogleDocsParser ──────────────────────────────────────────────────────────

class TestGoogleDocsParser:
    def _parse_with_mocks(self, source, opts=None):
        stubs, refs = _make_google_stubs()
        docx_bytes  = _make_docx_bytes()
        _make_drive_service(
            mime_type    = "application/vnd.google-apps.document",
            file_name    = "My Google Doc",
            export_bytes = docx_bytes,
            stubs        = refs,
        )
        opts = opts or ParseOptions()
        opts.extra["access_token"] = "ya29.token"

        with patch.dict(sys.modules, stubs):
            return GoogleDocsParser().parse(source, opts)

    def test_parses_from_url(self):
        doc = self._parse_with_mocks(DOCS_URL)
        assert doc is not None
        assert len(doc.sections) > 0

    def test_parses_from_bare_file_id(self):
        doc = self._parse_with_mocks(FAKE_FILE_ID)
        assert len(doc.sections) > 0

    def test_source_format_set_to_google_docs(self):
        doc = self._parse_with_mocks(DOCS_URL)
        assert doc.metadata.source_format == "google-docs"

    def test_source_path_preserved(self):
        doc = self._parse_with_mocks(DOCS_URL)
        assert doc.metadata.source_path == DOCS_URL

    def test_title_set_from_drive_file_name(self):
        doc = self._parse_with_mocks(DOCS_URL)
        assert doc.metadata.title == "My Google Doc"

    def test_raises_on_bytes_input(self):
        with pytest.raises(ParseError, match="does not accept raw bytes"):
            GoogleDocsParser().parse(b"binary content")

    def test_raises_on_invalid_file_id(self):
        opts = ParseOptions()
        opts.extra["access_token"] = "ya29.token"
        with pytest.raises(ParseError, match="Could not extract"):
            GoogleDocsParser().parse("not-a-valid-id", opts)

    def test_raises_on_missing_credentials(self, monkeypatch):
        stubs, refs = _make_google_stubs()
        monkeypatch.delenv("DISTILL_GOOGLE_CREDENTIALS", raising=False)

        with patch.dict(sys.modules, stubs):
            with pytest.raises(ParseError, match="requires credentials"):
                GoogleDocsParser().parse(DOCS_URL)

    def test_raises_on_wrong_mime_type(self):
        stubs, refs = _make_google_stubs()
        xlsx_bytes  = _make_xlsx_bytes()
        # Drive returns a Sheets file when we expect a Docs file
        _make_drive_service(
            mime_type    = "application/vnd.google-apps.spreadsheet",
            file_name    = "A Sheet",
            export_bytes = xlsx_bytes,
            stubs        = refs,
        )
        opts = ParseOptions()
        opts.extra["access_token"] = "ya29.token"

        with patch.dict(sys.modules, stubs):
            with pytest.raises(UnsupportedFormatError):
                GoogleDocsParser().parse(DOCS_URL, opts)

    def test_parses_from_gdoc_shortcut_path(self, tmp_path):
        p = tmp_path / f"{FAKE_FILE_ID}.gdoc"
        p.touch()
        doc = self._parse_with_mocks(str(p))
        assert len(doc.sections) > 0


# ── GoogleSheetsParser ────────────────────────────────────────────────────────

class TestGoogleSheetsParser:
    def _parse_with_mocks(self, source, opts=None):
        stubs, refs = _make_google_stubs()
        xlsx_bytes  = _make_xlsx_bytes()
        _make_drive_service(
            mime_type    = "application/vnd.google-apps.spreadsheet",
            file_name    = "My Google Sheet",
            export_bytes = xlsx_bytes,
            stubs        = refs,
        )
        opts = opts or ParseOptions()
        opts.extra["access_token"] = "ya29.token"

        with patch.dict(sys.modules, stubs):
            return GoogleSheetsParser().parse(source, opts)

    def test_parses_from_url(self):
        doc = self._parse_with_mocks(SHEETS_URL)
        assert doc is not None
        assert len(doc.sections) > 0

    def test_parses_from_bare_file_id(self):
        doc = self._parse_with_mocks(FAKE_FILE_ID)
        assert len(doc.sections) > 0

    def test_source_format_set_to_google_sheets(self):
        doc = self._parse_with_mocks(SHEETS_URL)
        assert doc.metadata.source_format == "google-sheets"

    def test_source_path_preserved(self):
        doc = self._parse_with_mocks(SHEETS_URL)
        assert doc.metadata.source_path == SHEETS_URL

    def test_title_set_from_drive_file_name(self):
        doc = self._parse_with_mocks(SHEETS_URL)
        assert doc.metadata.title == "My Google Sheet"

    def test_raises_on_bytes_input(self):
        with pytest.raises(ParseError, match="does not accept raw bytes"):
            GoogleSheetsParser().parse(b"binary content")

    def test_raises_on_missing_credentials(self, monkeypatch):
        stubs, refs = _make_google_stubs()
        monkeypatch.delenv("DISTILL_GOOGLE_CREDENTIALS", raising=False)

        with patch.dict(sys.modules, stubs):
            with pytest.raises(ParseError, match="requires credentials"):
                GoogleSheetsParser().parse(SHEETS_URL)

    def test_parses_from_gsheet_shortcut_path(self, tmp_path):
        p = tmp_path / f"{FAKE_FILE_ID}.gsheet"
        p.touch()
        doc = self._parse_with_mocks(str(p))
        assert len(doc.sections) > 0


# ── GoogleSlidesParser ────────────────────────────────────────────────────────

class TestGoogleSlidesParser:
    def _parse_with_mocks(self, source, opts=None):
        stubs, refs = _make_google_stubs()
        pptx_bytes  = _make_pptx_bytes()
        _make_drive_service(
            mime_type    = "application/vnd.google-apps.presentation",
            file_name    = "My Google Slides",
            export_bytes = pptx_bytes,
            stubs        = refs,
        )
        opts = opts or ParseOptions()
        opts.extra["access_token"] = "ya29.token"

        with patch.dict(sys.modules, stubs):
            return GoogleSlidesParser().parse(source, opts)

    def test_parses_from_url(self):
        doc = self._parse_with_mocks(SLIDES_URL)
        assert doc is not None
        assert len(doc.sections) > 0

    def test_parses_from_bare_file_id(self):
        doc = self._parse_with_mocks(FAKE_FILE_ID)
        assert len(doc.sections) > 0

    def test_source_format_set_to_google_slides(self):
        doc = self._parse_with_mocks(SLIDES_URL)
        assert doc.metadata.source_format == "google-slides"

    def test_source_path_preserved(self):
        doc = self._parse_with_mocks(SLIDES_URL)
        assert doc.metadata.source_path == SLIDES_URL

    def test_title_set_from_drive_file_name(self):
        doc = self._parse_with_mocks(SLIDES_URL)
        assert doc.metadata.title == "My Google Slides"

    def test_raises_on_bytes_input(self):
        with pytest.raises(ParseError, match="does not accept raw bytes"):
            GoogleSlidesParser().parse(b"binary content")

    def test_raises_on_missing_credentials(self, monkeypatch):
        stubs, refs = _make_google_stubs()
        monkeypatch.delenv("DISTILL_GOOGLE_CREDENTIALS", raising=False)

        with patch.dict(sys.modules, stubs):
            with pytest.raises(ParseError, match="requires credentials"):
                GoogleSlidesParser().parse(SLIDES_URL)

    def test_parses_from_gslides_shortcut_path(self, tmp_path):
        p = tmp_path / f"{FAKE_FILE_ID}.gslides"
        p.touch()
        doc = self._parse_with_mocks(str(p))
        assert len(doc.sections) > 0
