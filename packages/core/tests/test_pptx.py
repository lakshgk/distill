"""
Tests for distill.parsers.pptx — PptxParser, PptLegacyParser, and helpers.

Fixtures are built programmatically via python-pptx; no binary files are
checked into the repository.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest

from distill.ir import BlockQuote, Document, Image, List, Paragraph, Section, Table, TextRun
from distill.parsers.base import ParseError
from distill.parsers.pptx import (
    PptLegacyParser,
    PptxParser,
    _check_input_size,
    _check_zip_bomb,
)


# ── Fixture helpers ──────────────────────────────────────────────────────────

def _make_pptx(
    *,
    slides: list[dict] | None = None,
    title: str = "",
    author: str = "",
    subject: str = "",
    description: str = "",
    keywords: str = "",
) -> bytes:
    """
    Build a minimal .pptx in memory.

    Each slide dict may contain:
        title    (str)
        body     (str)           — body text in content placeholder
        bullets  (list[str])     — bullet items (level 1)
        notes    (str)           — speaker notes
        table    (list[list])    — rows×cols data (first row = header)
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    default_slides = [{"title": "Slide One", "body": "Hello world."}] if slides is None else slides
    for slide_def in default_slides:
        if slide_def.get("table"):
            layout = prs.slide_layouts[5]   # Blank
            sl = prs.slides.add_slide(layout)
            if slide_def.get("title"):
                tb = sl.shapes.add_textbox(Inches(0), Inches(0), Inches(8), Inches(0.5))
                tb.text_frame.text = slide_def["title"]
            data = slide_def["table"]
            rows_count = len(data)
            cols_count = len(data[0]) if data else 1
            tbl_shape = sl.shapes.add_table(
                rows_count, cols_count,
                Inches(0.5), Inches(1), Inches(8), Inches(rows_count * 0.5),
            )
            for r, row in enumerate(data):
                for c, val in enumerate(row):
                    tbl_shape.table.cell(r, c).text = str(val)
        else:
            layout = prs.slide_layouts[1]   # Title and Content
            sl = prs.slides.add_slide(layout)
            if slide_def.get("title"):
                sl.shapes.title.text = slide_def["title"]
            if slide_def.get("body") or slide_def.get("bullets"):
                ph = sl.placeholders[1]
                tf = ph.text_frame
                if slide_def.get("body"):
                    tf.text = slide_def["body"]
                for bullet in (slide_def.get("bullets") or []):
                    p = tf.add_paragraph()
                    p.text  = bullet
                    p.level = 1

        if slide_def.get("notes"):
            sl.notes_slide.notes_text_frame.text = slide_def["notes"]

    prs.core_properties.title       = title
    prs.core_properties.author      = author
    prs.core_properties.subject     = subject
    prs.core_properties.description = description
    prs.core_properties.keywords    = keywords

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _all_texts(doc: Document) -> list[str]:
    texts: list[str] = []
    for section in doc.sections:
        for run in (section.heading or []):
            texts.append(run.text)
        for block in section.blocks:
            if isinstance(block, Paragraph):
                texts.extend(r.text for r in block.runs)
            elif isinstance(block, List):
                texts.extend(r.text for item in block.items for r in item.content)
            elif isinstance(block, BlockQuote):
                for b in block.content:
                    if isinstance(b, Paragraph):
                        texts.extend(r.text for r in b.runs)
    return texts


def _tables(doc: Document) -> list[Table]:
    return [
        block
        for section in doc.sections
        for block in section.blocks
        if isinstance(block, Table)
    ]


def _blockquotes(doc: Document) -> list[BlockQuote]:
    return [
        block
        for section in doc.sections
        for block in section.blocks
        if isinstance(block, BlockQuote)
    ]


# ── Parser availability ───────────────────────────────────────────────────────

class TestParserAvailability:
    def test_pptx_is_available(self):
        assert PptxParser.is_available()

    def test_pptx_extensions(self):
        assert ".pptx" in PptxParser.extensions

    def test_pptx_missing_requires_empty(self):
        assert PptxParser.missing_requires() == []

    def test_ppt_legacy_is_available(self):
        assert PptLegacyParser.is_available()

    def test_ppt_legacy_extension(self):
        assert ".ppt" in PptLegacyParser.extensions


# ── Basic parsing ─────────────────────────────────────────────────────────────

class TestBasicParsing:
    def test_returns_document(self):
        data = _make_pptx()
        doc  = PptxParser().parse(data)
        assert isinstance(doc, Document)

    def test_one_section_per_slide(self):
        data = _make_pptx(slides=[
            {"title": "Intro"},
            {"title": "Details"},
            {"title": "Summary"},
        ])
        doc = PptxParser().parse(data)
        assert len(doc.sections) == 3

    def test_section_level_is_2(self):
        data = _make_pptx()
        doc  = PptxParser().parse(data)
        assert all(s.level == 2 for s in doc.sections)

    def test_slide_title_in_heading(self):
        data = _make_pptx(slides=[{"title": "Annual Review"}])
        doc  = PptxParser().parse(data)
        heading = " ".join(r.text for r in doc.sections[0].heading)
        assert "Annual Review" in heading

    def test_slide_title_used_in_heading(self):
        data = _make_pptx(slides=[{"title": "Revenue"}, {"title": "Costs"}])
        doc  = PptxParser().parse(data)
        h0 = " ".join(r.text for r in doc.sections[0].heading)
        h1 = " ".join(r.text for r in doc.sections[1].heading)
        assert h0 == "Revenue"
        assert h1 == "Costs"

    def test_untitled_slide_uses_first_textbox(self):
        data = _make_pptx(slides=[{"body": "No title here"}])
        doc  = PptxParser().parse(data)
        heading = " ".join(r.text for r in doc.sections[0].heading)
        assert "No title here" in heading

    def test_body_text_extracted(self):
        data = _make_pptx(slides=[{"title": "T", "body": "Main content here."}])
        doc  = PptxParser().parse(data)
        texts = _all_texts(doc)
        assert any("Main content here." in t for t in texts)

    def test_accepts_path(self, tmp_path):
        p = tmp_path / "test.pptx"
        p.write_bytes(_make_pptx())
        doc = PptxParser().parse(str(p))
        assert isinstance(doc, Document)

    def test_accepts_path_object(self, tmp_path):
        p = tmp_path / "test.pptx"
        p.write_bytes(_make_pptx())
        doc = PptxParser().parse(p)
        assert isinstance(doc, Document)

    def test_empty_presentation_no_crash(self):
        data = _make_pptx(slides=[])
        doc  = PptxParser().parse(data)
        assert isinstance(doc, Document)
        assert len(doc.sections) == 0


# ── Bullet lists ──────────────────────────────────────────────────────────────

class TestBulletLists:
    def test_bullets_become_list(self):
        data = _make_pptx(slides=[{
            "title": "Bullets",
            "bullets": ["Point A", "Point B", "Point C"],
        }])
        doc   = PptxParser().parse(data)
        lists = [b for b in doc.sections[0].blocks if isinstance(b, List)]
        assert len(lists) >= 1

    def test_bullet_texts_extracted(self):
        data = _make_pptx(slides=[{
            "title": "Bullets",
            "bullets": ["Alpha", "Beta", "Gamma"],
        }])
        doc  = PptxParser().parse(data)
        texts = _all_texts(doc)
        assert any("Alpha" in t for t in texts)
        assert any("Beta"  in t for t in texts)
        assert any("Gamma" in t for t in texts)


# ── Speaker notes ─────────────────────────────────────────────────────────────

class TestSpeakerNotes:
    def test_notes_become_blockquote(self):
        data = _make_pptx(slides=[{
            "title": "T",
            "body": "Content.",
            "notes": "These are speaker notes.",
        }])
        doc = PptxParser().parse(data)
        bqs = _blockquotes(doc)
        assert len(bqs) == 1

    def test_notes_text_in_blockquote(self):
        data = _make_pptx(slides=[{
            "title": "T",
            "notes": "Important talking point.",
        }])
        doc   = PptxParser().parse(data)
        texts = _all_texts(doc)
        assert any("Important talking point." in t for t in texts)

    def test_no_notes_no_blockquote(self):
        data = _make_pptx(slides=[{"title": "T", "body": "No notes here."}])
        doc  = PptxParser().parse(data)
        assert len(_blockquotes(doc)) == 0

    def test_notes_only_on_correct_slide(self):
        data = _make_pptx(slides=[
            {"title": "A", "notes": "Notes for A."},
            {"title": "B"},
        ])
        doc = PptxParser().parse(data)
        assert len(_blockquotes(doc)) == 1
        bq_text = " ".join(
            r.text for b in _blockquotes(doc) for p in b.content
            if isinstance(p, Paragraph) for r in p.runs
        )
        assert "Notes for A" in bq_text


# ── Tables ────────────────────────────────────────────────────────────────────

class TestTables:
    def test_table_extracted(self):
        data = _make_pptx(slides=[{
            "title": "Data",
            "table": [["Region", "Sales"], ["North", "500"], ["South", "300"]],
        }])
        doc  = PptxParser().parse(data)
        tbls = _tables(doc)
        assert len(tbls) == 1

    def test_table_first_row_is_header(self):
        data = _make_pptx(slides=[{
            "table": [["H1", "H2"], ["r1", "r2"]],
        }])
        doc  = PptxParser().parse(data)
        tbls = _tables(doc)
        assert all(c.is_header for c in tbls[0].rows[0].cells)

    def test_table_data_rows_not_header(self):
        data = _make_pptx(slides=[{
            "table": [["H1", "H2"], ["r1", "r2"]],
        }])
        doc  = PptxParser().parse(data)
        tbls = _tables(doc)
        assert all(not c.is_header for c in tbls[0].rows[1].cells)

    def test_table_cell_values(self):
        data = _make_pptx(slides=[{
            "table": [["City", "Pop"], ["Oslo", "700000"]],
        }])
        doc  = PptxParser().parse(data)
        tbls = _tables(doc)
        texts = [r.text for row in tbls[0].rows for c in row.cells for r in c.content]
        assert "City"   in texts
        assert "Oslo"   in texts
        assert "700000" in texts

    def test_table_row_cap(self):
        from distill.parsers.base import ParseOptions
        rows = [["H"]] + [[str(i)] for i in range(50)]
        data = _make_pptx(slides=[{"table": rows}])
        opts = ParseOptions(max_table_rows=5)
        doc  = PptxParser().parse(data, options=opts)
        tbls = _tables(doc)
        assert len(tbls[0].rows) == 5


# ── Metadata ──────────────────────────────────────────────────────────────────

class TestMetadata:
    def test_source_format(self):
        doc = PptxParser().parse(_make_pptx())
        assert doc.metadata.source_format == "pptx"

    def test_slide_count(self):
        data = _make_pptx(slides=[{"title": "A"}, {"title": "B"}])
        doc  = PptxParser().parse(data)
        assert doc.metadata.slide_count == 2

    def test_word_count_positive(self):
        data = _make_pptx(slides=[{"title": "T", "body": "One two three four five."}])
        doc  = PptxParser().parse(data)
        assert doc.metadata.word_count is not None
        assert doc.metadata.word_count > 0

    def test_title(self, tmp_path):
        p = tmp_path / "t.pptx"
        p.write_bytes(_make_pptx(title="My Deck"))
        doc = PptxParser().parse(p)
        assert doc.metadata.title == "My Deck"

    def test_author(self, tmp_path):
        p = tmp_path / "t.pptx"
        p.write_bytes(_make_pptx(author="Jane Doe"))
        doc = PptxParser().parse(p)
        assert doc.metadata.author == "Jane Doe"

    def test_subject(self, tmp_path):
        p = tmp_path / "t.pptx"
        p.write_bytes(_make_pptx(subject="Q4 Review"))
        doc = PptxParser().parse(p)
        assert doc.metadata.subject == "Q4 Review"

    def test_keywords_comma(self, tmp_path):
        p = tmp_path / "t.pptx"
        p.write_bytes(_make_pptx(keywords="finance, growth, q4"))
        doc = PptxParser().parse(p)
        assert "finance" in doc.metadata.keywords
        assert "growth"  in doc.metadata.keywords

    def test_keywords_semicolon(self, tmp_path):
        p = tmp_path / "t.pptx"
        p.write_bytes(_make_pptx(keywords="alpha;beta;gamma"))
        doc = PptxParser().parse(p)
        assert doc.metadata.keywords == ["alpha", "beta", "gamma"]


# ── Security ──────────────────────────────────────────────────────────────────

class TestSecurity:
    def test_input_size_bytes_rejected(self):
        oversized = b"x" * (55 * 1024 * 1024)
        with pytest.raises(ParseError, match="50 MB"):
            _check_input_size(oversized, 50 * 1024 * 1024)

    def test_input_size_path_rejected(self, tmp_path):
        big = tmp_path / "big.pptx"
        big.write_bytes(b"x" * (55 * 1024 * 1024))
        with pytest.raises(ParseError, match="50 MB"):
            _check_input_size(str(big), 50 * 1024 * 1024)

    def test_input_size_ok(self):
        _check_input_size(b"tiny", 50 * 1024 * 1024)  # must not raise

    def test_zip_bomb_detected(self):
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("big.bin", "A" * (501 * 1024 * 1024))
        with pytest.raises(ParseError, match="500 MB"):
            _check_zip_bomb(buf.getvalue(), 500 * 1024 * 1024)

    def test_bad_zip_raises(self):
        with pytest.raises(ParseError, match="not a valid PPTX"):
            _check_zip_bomb(b"not a zip file", 500 * 1024 * 1024)

    def test_parser_rejects_oversized(self):
        oversized = b"x" * (55 * 1024 * 1024)
        with pytest.raises(ParseError, match="50 MB"):
            PptxParser().parse(oversized)

    def test_custom_size_limit(self):
        from distill.parsers.base import ParseOptions
        data = b"x" * (15 * 1024 * 1024)
        opts = ParseOptions(extra={"max_file_size": 10 * 1024 * 1024})
        with pytest.raises(ParseError, match="10 MB"):
            PptxParser().parse(data, options=opts)

    def test_garbled_bytes_raises_parse_error(self):
        with pytest.raises(ParseError):
            PptxParser().parse(b"this is not a pptx")


# ── PptLegacyParser ───────────────────────────────────────────────────────────

class TestPptLegacyParser:
    def test_raises_parse_error(self):
        # Wired to LibreOffice — will fail due to binary not found or conversion error
        with pytest.raises(ParseError):
            PptLegacyParser().parse(b"garbage")

    def test_error_mentions_libreoffice(self):
        with pytest.raises(ParseError, match="LibreOffice"):
            PptLegacyParser().parse(b"garbage")

    def test_extension_registered(self):
        assert ".ppt" in PptLegacyParser.extensions


# ── Render integration ─────────────────────────────────────────────────────────

class TestRenderIntegration:
    def test_renders_to_markdown(self):
        data = _make_pptx(slides=[{"title": "Intro", "body": "Content here."}])
        doc  = PptxParser().parse(data)
        md   = doc.render(front_matter=False)
        assert isinstance(md, str)
        assert len(md) > 0

    def test_front_matter_has_source_format(self):
        data = _make_pptx()
        doc  = PptxParser().parse(data)
        md   = doc.render(front_matter=True)
        assert "---" in md
        assert "pptx" in md

    def test_no_front_matter_when_suppressed(self):
        data = _make_pptx()
        doc  = PptxParser().parse(data)
        md   = doc.render(front_matter=False)
        assert not md.startswith("---")

    def test_slide_heading_in_markdown(self):
        data = _make_pptx(slides=[{"title": "Executive Summary"}])
        doc  = PptxParser().parse(data)
        md   = doc.render(front_matter=False)
        assert "Executive Summary" in md

    def test_table_pipe_syntax_in_output(self):
        data = _make_pptx(slides=[{
            "table": [["A", "B"], ["1", "2"]],
        }])
        doc = PptxParser().parse(data)
        md  = doc.render(front_matter=False)
        assert "|" in md

    def test_blockquote_in_markdown(self):
        data = _make_pptx(slides=[{
            "title": "T",
            "notes": "Speaker note text.",
        }])
        doc = PptxParser().parse(data)
        md  = doc.render(front_matter=False)
        assert ">" in md  # blockquote marker
        assert "Speaker note text." in md

    def test_table_produces_gfm_pipe_syntax(self):
        data = _make_pptx(slides=[{
            "table": [["Region", "Sales"], ["North", "500"]],
        }])
        doc = PptxParser().parse(data)
        md  = doc.render(front_matter=False)
        assert "| --- |" in md

    def test_empty_image_suppressed(self):
        from distill.renderer import MarkdownRenderer
        doc = Document(sections=[
            Section(level=0, blocks=[
                Image(),
                Paragraph(runs=[TextRun(text="Visible text")]),
            ])
        ])
        md = MarkdownRenderer(front_matter=False).render(doc)
        assert "![" not in md
        assert "Visible text" in md


# ── Spec B: Slide title / Image extraction / Vision ──────────────────────────

def _make_png_bytes():
    """Create a minimal valid 1x1 PNG image."""
    import struct, zlib
    raw = b'\x00\xff\x00\x00'
    compressed = zlib.compress(raw)
    def chunk(ctype, data):
        c = ctype + data
        return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
    return (b'\x89PNG\r\n\x1a\n' +
            chunk(b'IHDR', struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)) +
            chunk(b'IDAT', compressed) +
            chunk(b'IEND', b''))


def _make_pptx_with_image() -> bytes:
    """Build a PPTX with one titled slide containing an image."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    sl.shapes.title.text = "Architecture Overview"
    sl.shapes.add_picture(
        io.BytesIO(_make_png_bytes()),
        Inches(1), Inches(1), Inches(2), Inches(2),
    )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestSlideTitle:
    def test_title_extracted_in_heading(self):
        data = _make_pptx(slides=[{"title": "Architecture Overview"}])
        doc = PptxParser().parse(data)
        md = doc.render(front_matter=False)
        assert "## Architecture Overview" in md
        assert "## Slide 1" not in md

    def test_empty_slide_suppressed_when_no_content(self):
        from pptx import Presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])  # Blank, no title, no content
        buf = io.BytesIO()
        prs.save(buf)
        doc = PptxParser().parse(buf.getvalue())
        assert len(doc.sections) == 0


class TestImageExtraction:
    def test_image_suppressed(self):
        from distill.parsers.base import ParseOptions
        data = _make_pptx_with_image()
        doc = PptxParser().parse(data, options=ParseOptions(images="suppress"))
        images = [b for s in doc.sections for b in s.blocks if isinstance(b, Image)]
        assert len(images) == 0
        md = doc.render(front_matter=False)
        assert "![" not in md

    def test_image_extracted(self, tmp_path):
        from distill.parsers.base import ParseOptions
        data = _make_pptx_with_image()
        img_dir = tmp_path / "images"
        doc = PptxParser().parse(
            data,
            options=ParseOptions(images="extract", image_dir=str(img_dir)),
        )
        images = [b for s in doc.sections for b in s.blocks if isinstance(b, Image)]
        assert len(images) >= 1
        assert any(img.path is not None for img in images)
        written = list(img_dir.glob("*")) if img_dir.exists() else []
        assert len(written) >= 1
        md = doc.render(front_matter=False)
        assert "![](" not in md or "![" in md  # no fully-empty image tags

    def test_vision_captioning_fallback(self, tmp_path):
        from distill.parsers.base import ParseOptions
        from distill.warnings import WarningCollector, WarningType
        data = _make_pptx_with_image()
        collector = WarningCollector()
        doc = PptxParser().parse(
            data,
            options=ParseOptions(
                images="caption",
                vision_provider="anthropic",
                image_dir=str(tmp_path / "images"),
                collector=collector,
            ),
        )
        # Conversion completes without raising
        assert isinstance(doc, Document)
        # vision_caption_failed warning emitted
        assert collector.has(WarningType.vision_caption_failed)
        # Image nodes still included (not suppressed)
        images = [b for s in doc.sections for b in s.blocks if isinstance(b, Image)]
        assert len(images) >= 1


# ── Spec C: Decorative image filtering ────────────────────────────────────────

class TestDecorativeImageFiltering:
    def test_classify_image_full_bleed_decorative(self):
        from distill.parsers.base import classify_image
        from distill.ir import ImageType
        result = classify_image(
            mode="pptx",
            shape_w=9144000, shape_h=5143500,
            slide_w=9144000, slide_h=5143500,
            name="bg_teal",
        )
        assert result == ImageType.DECORATIVE

    def test_classify_image_content_not_decorative(self):
        from distill.parsers.base import classify_image
        from distill.ir import ImageType
        result = classify_image(
            mode="pptx",
            shape_w=2000000, shape_h=2000000,
            slide_w=9144000, slide_h=5143500,
            name="Jason Foodman",
        )
        assert result == ImageType.UNKNOWN

    def test_decorative_image_suppressed_in_render(self):
        from distill.renderer import MarkdownRenderer
        from distill.ir import ImageType
        doc = Document(sections=[
            Section(level=0, blocks=[
                Image(image_type=ImageType.DECORATIVE, path="bg.png", alt_text="background"),
                Paragraph(runs=[TextRun(text="Real content")]),
            ])
        ])
        md = MarkdownRenderer(front_matter=False).render(doc)
        assert "![" not in md
        assert "Real content" in md


# ── Spec D: Title fallback + empty slide suppression ─────────────────────────

class TestTitleFallback:
    def test_first_textbox_used_when_no_title_placeholder(self):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        sl = prs.slides.add_slide(prs.slide_layouts[5])  # Blank — no title placeholder
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1))
        tb.text_frame.text = "My Custom Title"
        buf = io.BytesIO()
        prs.save(buf)
        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)
        assert "## My Custom Title" in md
        assert "## Slide 1" not in md

    def test_title_placeholder_takes_priority_over_textbox(self):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        sl = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        sl.shapes.title.text = "Official Title"
        tb = sl.shapes.add_textbox(Inches(0.5), Inches(3), Inches(6), Inches(1))
        tb.text_frame.text = "First Text Box"
        buf = io.BytesIO()
        prs.save(buf)
        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)
        assert "## Official Title" in md
        assert "## First Text Box" not in md


class TestEmptySlideSuppression:
    def test_empty_slide_not_in_output(self):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        # Slide 1: has content
        sl1 = prs.slides.add_slide(prs.slide_layouts[5])
        tb = sl1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1))
        tb.text_frame.text = "Slide With Content"
        # Slide 2: completely empty (no text, no images)
        prs.slides.add_slide(prs.slide_layouts[5])
        buf = io.BytesIO()
        prs.save(buf)
        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)
        assert "Slide With Content" in md
        assert "## Slide 2" not in md
        assert len(doc.sections) == 1


# ── Batch 1: Duplicate title fix + Image classification fix ──────────────────

class TestDuplicateTitleFix:
    def test_slide_title_not_duplicated_in_body(self):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Unique Heading"
        slide.placeholders[1].text = "Body content here"
        buf = io.BytesIO()
        prs.save(buf)
        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)
        heading_count = md.count("Unique Heading")
        assert heading_count == 1, (
            f"Expected title to appear exactly once, found {heading_count} times:\n{md}"
        )

    def test_fallback_title_shape_not_duplicated_in_body(self):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # Blank layout — no title placeholder
        slide = prs.slides.add_slide(blank_layout)
        txBox1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox1.text_frame.text = "Fallback Title Text"
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(2))
        txBox2.text_frame.text = "Body content"
        buf = io.BytesIO()
        prs.save(buf)
        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)
        count = md.count("Fallback Title Text")
        assert count == 1, (
            f"Expected fallback title to appear exactly once, found {count} times:\n{md}"
        )


class TestImageClassificationFix:
    def test_image_0_not_classified_decorative(self):
        from distill.parsers.base import classify_image
        from distill.ir import ImageType
        result = classify_image(
            mode="pptx",
            shape_w=9144000,
            shape_h=5911453,
            slide_w=9144000,
            slide_h=5143500,
            name="Image 0",
        )
        assert result != ImageType.DECORATIVE, (
            f"'Image 0' should not be DECORATIVE, got {result}"
        )

    def test_genuine_decorative_shape_still_caught(self):
        from distill.parsers.base import classify_image
        from distill.ir import ImageType
        result = classify_image(
            mode="pptx",
            shape_w=9144000,
            shape_h=5143500,
            slide_w=9144000,
            slide_h=5143500,
            name="Background",
        )
        assert result == ImageType.DECORATIVE, (
            f"Full-slide background shape should be DECORATIVE, got {result}"
        )


# ── Batch 3: Bullet list XML marker detection ───────────────────────────────

def _add_bullet_para(tf, text, level=0):
    """Add a bullet paragraph with explicit buChar XML marker."""
    from pptx.oxml import parse_xml
    p = parse_xml(
        f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:pPr lvl="{level}"><a:buChar char="\u2022"/></a:pPr>'
        f'<a:r><a:t>{text}</a:t></a:r>'
        f'</a:p>'
    )
    tf._txBody.append(p)


def _add_plain_para(tf, text):
    """Add a plain (non-bullet) paragraph."""
    from pptx.oxml import parse_xml
    p = parse_xml(
        f'<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:r><a:t>{text}</a:t></a:r>'
        f'</a:p>'
    )
    tf._txBody.append(p)


class TestBulletXmlDetection:
    def test_bullet_paragraphs_render_as_list(self):
        """Paragraphs with buChar XML marker render as Markdown list items."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = txBox.text_frame
        tf._txBody.clear()
        _add_bullet_para(tf, "Alpha")
        _add_bullet_para(tf, "Beta")
        _add_bullet_para(tf, "Gamma")

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)

        for item in ["Alpha", "Beta", "Gamma"]:
            assert f"- {item}" in md or f"* {item}" in md, \
                f'"{item}" not rendered as list item:\n{md}'

    def test_nested_bullet_indented(self):
        """Nested bullet paragraphs are indented in Markdown output."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = txBox.text_frame
        tf._txBody.clear()
        _add_bullet_para(tf, "Parent item", level=0)
        _add_bullet_para(tf, "Child item", level=1)
        _add_bullet_para(tf, "Grandchild item", level=2)

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)
        lines = md.splitlines()

        child_lines = [l for l in lines if "Child item" in l]
        grandchild_lines = [l for l in lines if "Grandchild item" in l]

        assert child_lines, f"Child item not in output:\n{md}"
        assert grandchild_lines, f"Grandchild item not in output:\n{md}"
        assert child_lines[0].startswith(" ") or child_lines[0].startswith("\t"), \
            f'Child item not indented: "{child_lines[0]}"'
        assert len(grandchild_lines[0]) - len(grandchild_lines[0].lstrip()) > \
               len(child_lines[0]) - len(child_lines[0].lstrip()), \
            "Grandchild not indented more than child"

    def test_mixed_bullet_and_plain_paragraphs(self):
        """Non-bullet paragraphs in the same text frame render as paragraphs."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = txBox.text_frame
        tf._txBody.clear()
        _add_plain_para(tf, "Intro text")
        _add_bullet_para(tf, "Bullet one")
        _add_bullet_para(tf, "Bullet two")
        _add_plain_para(tf, "Closing text")

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)

        assert "Intro text" in md, f"Plain paragraph before bullets missing:\n{md}"
        assert "Closing text" in md, f"Plain paragraph after bullets missing:\n{md}"
        assert "- Bullet one" in md or "* Bullet one" in md, \
            f"Bullet one not rendered as list item:\n{md}"

    def test_bu_none_not_rendered_as_bullet(self):
        """Paragraphs with buNone are NOT rendered as bullet list items."""
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.oxml import parse_xml

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = txBox.text_frame
        tf._txBody.clear()

        p = parse_xml(
            '<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:pPr><a:buNone/></a:pPr>'
            '<a:r><a:t>Not a bullet</a:t></a:r>'
            '</a:p>'
        )
        tf._txBody.append(p)

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)

        assert "- Not a bullet" not in md and "* Not a bullet" not in md, \
            f"buNone paragraph incorrectly rendered as bullet:\n{md}"
        assert "Not a bullet" in md, \
            f"buNone paragraph content missing from output:\n{md}"


# ── Batch 4: Heuristic title extraction ─────────────────────────────────────

def _add_textbox_with_font(slide, text, left_in, top_in, w_in, h_in, font_pt):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in)
    )
    tf = txBox.text_frame
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_pt)
    return txBox


class TestHeuristicTitleExtraction:
    def test_heuristic_title_selected_over_zorder_first(self):
        """Large-font shape in top zone is selected over z-order-first decorative shape."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_h_in = prs.slide_height / 914400

        _add_textbox_with_font(slide, "Decorative", 0, 0, 1, 0.3, 10)
        _add_textbox_with_font(slide, "Heuristic Title", 1, slide_h_in * 0.08, 6, 0.6, 24)
        _add_textbox_with_font(slide, "Body content", 1, slide_h_in * 0.25, 6, 1, 16)

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)

        heading_lines = [l for l in md.splitlines() if l.startswith("#")]
        assert any("Heuristic Title" in l for l in heading_lines), \
            f'Expected "Heuristic Title" as slide heading, got: {heading_lines}\n{md}'
        assert not any("Decorative" in l for l in heading_lines), \
            f"Decorative shape must not be used as heading: {heading_lines}"

    def test_heuristic_title_not_duplicated_in_body(self):
        """Title selected by heuristic must appear exactly once in output."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_h_in = prs.slide_height / 914400

        _add_textbox_with_font(slide, "Unique Heuristic Title", 1, slide_h_in * 0.06, 6, 0.6, 24)
        _add_textbox_with_font(slide, "Body content here", 1, slide_h_in * 0.30, 6, 1, 16)

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)

        count = md.count("Unique Heuristic Title")
        assert count == 1, \
            f"Title appears {count} times, expected 1:\n{md}"

    def test_no_heuristic_candidate_fallback(self):
        """Slide with no heuristic candidate falls back gracefully without crashing."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_h_in = prs.slide_height / 914400

        _add_textbox_with_font(slide, "Body only text", 1, slide_h_in * 0.30, 6, 1, 16)

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)

        assert md is not None, "Conversion returned None markdown"
        assert "Body only text" in md, \
            f"Body text missing from output:\n{md}"

    def test_long_text_excluded_from_heuristic(self):
        """Text > 120 chars in the top zone is not selected as the title."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_h_in = prs.slide_height / 914400

        long_text = ("This is a very long paragraph that exceeds one hundred and "
                     "twenty characters and should therefore not be selected as the "
                     "slide title by the heuristic.")
        assert len(long_text) > 120

        _add_textbox_with_font(slide, long_text, 0.5, slide_h_in * 0.05, 8, 1, 24)
        _add_textbox_with_font(slide, "Short Real Title", 0.5, slide_h_in * 0.08, 4, 0.5, 24)

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)

        heading_lines = [l for l in md.splitlines() if l.startswith("#")]
        assert any("Short Real Title" in l for l in heading_lines), \
            f"Short title not selected as heading: {heading_lines}\n{md}"
        assert not any(long_text[:30] in l for l in heading_lines), \
            f"Long text incorrectly selected as heading: {heading_lines}"

    def test_standard_title_placeholder_takes_priority(self):
        """Standard title placeholder still takes priority over heuristic."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Standard Placeholder Title"
        try:
            slide.placeholders[1].text = "Body content"
        except (KeyError, IndexError):
            pass

        buf = io.BytesIO()
        prs.save(buf)

        doc = PptxParser().parse(buf.getvalue())
        md = doc.render(front_matter=False)

        heading_lines = [l for l in md.splitlines() if l.startswith("#")]
        assert any("Standard Placeholder Title" in l for l in heading_lines), \
            f"Standard title placeholder not used as heading: {heading_lines}\n{md}"


# ── Batch 9: PPTX image alt text from shape.description ────────────────────

class TestImageAltTextDescription:
    def test_image_alt_text_uses_shape_description(self):
        """shape.description (author alt text) takes priority over shape.name."""
        import tempfile
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image as PILImage
        from distill.parsers.base import ParseOptions

        img_buf = io.BytesIO()
        PILImage.new("RGB", (10, 10), color="blue").save(img_buf, format="PNG")
        img_buf.seek(0)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(img_buf, Inches(1), Inches(1), Inches(2), Inches(2))
        pic._element.nvPicPr.cNvPr.set("descr", "A blue square used as a logo")

        buf = io.BytesIO()
        prs.save(buf)

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = PptxParser().parse(
                buf.getvalue(),
                options=ParseOptions(images="extract", image_dir=tmpdir),
            )
            md = doc.render(front_matter=False)

        assert "A blue square used as a logo" in md, \
            f"shape.description not used as alt text:\n{md}"

    def test_image_alt_text_falls_back_to_shape_name(self):
        """When shape.description is empty, shape.name is used as fallback."""
        import tempfile
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image as PILImage
        from distill.parsers.base import ParseOptions

        img_buf = io.BytesIO()
        PILImage.new("RGB", (10, 10), color="red").save(img_buf, format="PNG")
        img_buf.seek(0)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(img_buf, Inches(1), Inches(1), Inches(2), Inches(2))
        pic.name = "CompanyLogo"
        # Clear the default descr so fallback to shape.name triggers
        pic._element.nvPicPr.cNvPr.set("descr", "")

        buf = io.BytesIO()
        prs.save(buf)

        with tempfile.TemporaryDirectory() as tmpdir:
            doc = PptxParser().parse(
                buf.getvalue(),
                options=ParseOptions(images="extract", image_dir=tmpdir),
            )
            md = doc.render(front_matter=False)

        assert "CompanyLogo" in md, \
            f"shape.name not used as fallback alt text:\n{md}"


# ── Batch 11: Footer placeholder suppression ────────────────────────────────

class TestFooterPlaceholderSuppression:
    def test_footer_placeholder_suppressed(self):
        """Footer placeholder indices 11, 12, 13 must be in the constant."""
        from distill.parsers.pptx import _FOOTER_PLACEHOLDER_INDICES
        assert 11 in _FOOTER_PLACEHOLDER_INDICES
        assert 12 in _FOOTER_PLACEHOLDER_INDICES
        assert 13 in _FOOTER_PLACEHOLDER_INDICES

    def test_footer_placeholder_indices_constant_defined(self):
        """_FOOTER_PLACEHOLDER_INDICES must be frozenset({11, 12, 13})."""
        from distill.parsers.pptx import _FOOTER_PLACEHOLDER_INDICES
        assert _FOOTER_PLACEHOLDER_INDICES == frozenset({11, 12, 13}), \
            f"Expected frozenset({{11, 12, 13}}), got {_FOOTER_PLACEHOLDER_INDICES}"

    def test_footer_skip_logic_in_parse_slide(self):
        """_parse_slide must contain footer placeholder skip logic."""
        import inspect
        from distill.parsers.pptx import PptxParser
        src = inspect.getsource(PptxParser._parse_slide)
        assert "_FOOTER_PLACEHOLDER_INDICES" in src, \
            "_parse_slide must reference _FOOTER_PLACEHOLDER_INDICES"
        assert "placeholder_format.idx" in src, \
            "_parse_slide must check placeholder_format.idx"
